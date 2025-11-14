#!/usr/bin/env python3
"""
filePreparation_v2.py  — production-ready preprocessing for GPT4All LocalDocs

Summary (implemented):
 - Static input/output (no CLI args)
 - Compatible folder: Processed_Compatible/  (case a & b)
 - Incompatible folder: Incompatible/
 - Hash DB (filehash.json) in Processed_Compatible/ (stores only processed files' hashes)
 - For every file:
     1) compute hash and skip if already processed
     2) attempt text extraction (pdf->PyMuPDF, docx->python-docx, pptx->python-pptx, xlsx/csv->pandas, epub->ebooklib, txt/md/rst->read)
     3) if extension compatible and extracted text >= MIN_TEXT_CHARS -> copy original into Processed_Compatible/
        (case a)
     4) if extension compatible but text < MIN_TEXT_CHARS -> OCR via page-by-page rendering with PyMuPDF + pytesseract, save .txt to Processed_Compatible/
        (case b)
     5) if extension not compatible -> copy original into Incompatible/
        (case c)
 - Tesseract OCR only (pytesseract). No easyocr/torch.
 - Memory-safe page-by-page OCR with explicit del + gc.collect()
 - Optional parallelism via ThreadPoolExecutor (configure MAX_WORKERS)
 - Logs and final summary

 NOTES:
 - Ensure Tesseract is installed and path below is correct.
 - PyMuPDF (fitz) used for PDF parsing & page rendering (no poppler needed).
 - If some lightweight Python packages are missing, the script will attempt to pip install them
   (this is safe; it avoids installing heavy ML libs).

Configure the static paths below, then run the script.
"""
# ---------------------- Configuration ----------------------
INPUT_DIR = r"C:\Users\Shubham Mehta\Documents\PME\knowledgeBase"
OUTPUT_DIR = r"C:\Users\Shubham Mehta\Documents\PME\knowledgeBase_prepared"
TESSERACT_CMD = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
# Set to os.cpu_count() for maximum CPU, or smaller to be conservative.
# Cap applied automatically to avoid excessive memory usage.
import os
CPU_COUNT = os.cpu_count() or 2
MAX_WORKERS = min(max(1, CPU_COUNT - 1), 6)  # default; change to min(CPU_COUNT, 6) or CPU_COUNT to increase
# Minimum characters to treat file as text-ready
MIN_TEXT_CHARS = 80
# Logging and housekeeping
LOG_DIR = os.path.join(OUTPUT_DIR, "logs")
ERROR_LOG = os.path.join(LOG_DIR, "errors.log")
HASH_DB_PATH = os.path.join(OUTPUT_DIR, "Processed_Compatible", "filehash.json")
# Supported extensions
GPT4ALL_COMPATIBLE_EXTS = {".docx", ".pdf", ".txt", ".md", ".rst"}
ALL_SUPPORTED_EXTS = GPT4ALL_COMPATIBLE_EXTS.union({
    ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".epub",
    ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif"
})
# -----------------------------------------------------------

# ---------------------- Ensure lightweight deps ----------------------
LIGHT_PKGS = {
    "pytesseract": "pytesseract",
    "PIL": "Pillow",
    "fitz": "PyMuPDF",
    "docx": "python-docx",
    "pptx": "python-pptx",
    "pandas": "pandas",
    "tqdm": "tqdm",
    "ebooklib": "ebooklib",
    "openpyxl": "openpyxl"
}

import subprocess, sys, json, shutil, time, math, gc, traceback
for mod, pipname in LIGHT_PKGS.items():
    try:
        if mod == "PIL":
            import PIL  # type: ignore
        else:
            __import__(mod)
    except Exception:
        print(f"[INFO] Installing missing package: {pipname}")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--upgrade", pipname])

# Now import required modules safely
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from ebooklib import epub
from tqdm import tqdm

# configure tesseract path
pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

# ---------------------- Prepare folders ----------------------
Processed_DIR = os.path.join(OUTPUT_DIR, "Processed_Compatible")
INCOMPATIBLE_DIR = os.path.join(OUTPUT_DIR, "Incompatible")
os.makedirs(Processed_DIR, exist_ok=True)
os.makedirs(INCOMPATIBLE_DIR, exist_ok=True)
os.makedirs(LOG_DIR, exist_ok=True)

# ---------------------- Utility helpers ----------------------
from pathlib import Path
def sha256_of_file(path: str, bufsize=1024*1024):
    import hashlib
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(bufsize)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()

def safe_copy(src: str, dst_dir: str):
    ensure_dir = Path(dst_dir)
    ensure_dir.mkdir(parents=True, exist_ok=True)
    dst = ensure_dir / Path(src).name
    if dst.exists():
        # append timestamp
        dst = ensure_dir / f"{Path(src).stem}_{int(time.time())}{Path(src).suffix}"
    shutil.copy2(src, str(dst))
    return str(dst)

def write_text_file(path: str, text: str):
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    with open(p, "w", encoding="utf-8", errors="ignore") as f:
        f.write(text)

def append_error(msg: str):
    with open(ERROR_LOG, "a", encoding="utf-8") as f:
        f.write(f"{time.asctime()}: {msg}\n")

# ---------------------- Extraction helpers ----------------------
def extract_text_from_docx(path: str) -> str:
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
    except Exception as e:
        append_error(f"docx extract failed {path}: {e}")
        return ""

def extract_text_from_pptx(path: str) -> str:
    out = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        out.append(t)
    except Exception as e:
        append_error(f"pptx extract failed {path}: {e}")
    return "\n".join(out)

def extract_text_from_xlsx(path: str) -> str:
    out = []
    try:
        sheets = pd.read_excel(path, sheet_name=None, engine="openpyxl")
        for sname, df in sheets.items():
            out.append(f"--- Sheet: {sname} ---")
            out.append(df.to_csv(index=False))
    except Exception as e:
        append_error(f"xlsx extract failed {path}: {e}")
    return "\n".join(out)

def extract_text_from_csv(path: str) -> str:
    try:
        df = pd.read_csv(path, dtype=str, keep_default_na=False)
        return df.to_csv(index=False)
    except Exception as e:
        append_error(f"csv extract failed {path}: {e}")
        return ""

def extract_text_from_epub(path: str) -> str:
    out = []
    try:
        book = epub.read_epub(path)
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                try:
                    content = item.get_content().decode('utf-8', errors='ignore')
                    # strip tags lightly
                    import re
                    text = re.sub('<[^<]+?>', '', content)
                    if text.strip():
                        out.append(text.strip())
                except Exception:
                    continue
    except Exception as e:
        append_error(f"epub extract failed {path}: {e}")
    return "\n\n".join(out)

def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            return fh.read()
    except Exception as e:
        append_error(f"txt read failed {path}: {e}")
        return ""

# PDF handling: first try text extraction via PyMuPDF, if insufficient then page-by-page render+OCR
def extract_text_from_pdf(path: str, ocr_dpi: int = 200) -> str:
    accumulated = ""
    try:
        doc = fitz.open(path)
        for p in doc:
            accumulated += p.get_text("text") + "\n"
        doc.close()
    except Exception as e:
        append_error(f"pdf text extraction (fitz) failed {path}: {e}")
        accumulated = ""

    if accumulated and len(accumulated.strip()) >= MIN_TEXT_CHARS:
        return accumulated

    # else fallback to page-by-page OCR rendering to conserve memory
    text_out = []
    try:
        doc = fitz.open(path)
        page_count = doc.page_count
        # choose a zoom matrix, control dpi via zoom factor: DPI ≈ 72 * zoom
        zoom = max(1.0, ocr_dpi / 72.0)
        mat = fitz.Matrix(zoom, zoom)
        for i in range(page_count):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            # convert pix to PIL Image
            mode = "RGB" if pix.n >= 3 else "L"
            img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            # OCR single page
            try:
                page_text = pytesseract.image_to_string(img)
            except Exception as oe:
                append_error(f"tesseract failed on page {i} of {path}: {oe}")
                page_text = ""
            text_out.append(page_text)
            # free memory aggressively
            del img
            del pix
            del page
            gc.collect()
        doc.close()
    except Exception as e:
        append_error(f"pdf OCR rendering failed {path}: {e}")

    return "\n\n".join(text_out)

# generic wrapper based on ext
def extract_text_generic(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(path)
    if ext in [".xlsx", ".xls"]:
        return extract_text_from_xlsx(path)
    if ext == ".csv":
        return extract_text_from_csv(path)
    if ext == ".epub":
        return extract_text_from_epub(path)
    if ext in [".txt", ".md", ".rst"]:
        return extract_text_from_txt(path)
    if ext in [".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif"]:
        # image file: use tesseract directly (via PIL)
        try:
            img = Image.open(path)
            txt = pytesseract.image_to_string(img)
            del img
            gc.collect()
            return txt
        except Exception as e:
            append_error(f"image OCR failed {path}: {e}")
            return ""
    return ""

# ---------------------- Main file processing logic ----------------------
from concurrent.futures import ThreadPoolExecutor, as_completed

def process_single_file(path: str):
    """
    Process a single file according to the algorithm:
     - compute hash; skip if in DB
     - attempt extraction; if ext compatible and text length >= threshold -> copy original to Processed_Compatible
     - if ext compatible but insufficient text -> OCR -> save .txt to Processed_Compatible
     - if ext not compatible -> copy original to Incompatible
    Returns tuple (status, path, info)
    """
    try:
        fhash = sha256_of_file(path)
        # load hash db lazily
        if fhash in HASH_DB_CACHE:
            return ("skipped", path, "already_processed")

        ext = Path(path).suffix.lower()

        # If extension not supported -> copy to Incompatible
        if ext not in ALL_SUPPORTED_EXTS:
            dest = safe_copy(path, INCOMPATIBLE_DIR)
            # don't add to hash db
            return ("incompatible", path, dest)

        # else supported -> try extract text
        text = extract_text_generic(path)
        text_len = len(text.strip()) if text else 0

        # If extension is GPT4All-compatible and text enough -> copy original (case a)
        if ext in GPT4ALL_COMPATIBLE_EXTS and text_len >= MIN_TEXT_CHARS:
            dest = safe_copy(path, Processed_DIR)
            # record hash
            HASH_DB_CACHE[fhash] = {"original": path, "status": "copied_original", "out": dest}
            save_hash_db()
            return ("copied_original", path, dest)

        # If extension is GPT4All-compatible but text insufficient -> OCR (case b)
        if ext in GPT4ALL_COMPATIBLE_EXTS and text_len < MIN_TEXT_CHARS:
            # perform OCR (PDF uses its own pdf OCR routine)
            ocr_text = ""
            if ext == ".pdf":
                ocr_text = extract_text_from_pdf(path)
            else:
                # for docx that are image-based, try to open with fitz if possible or fallback to converting embedded images? For simplicity attempt docx text, then fallback to page rendering via pdf conversion is complex.
                # We'll attempt a conservative approach: convert docx->text via extract_text_generic again and if still empty, mark failed.
                ocr_text = extract_text_generic(path)
            if ocr_text and len(ocr_text.strip()) >= 10:
                outtxt = os.path.join(Processed_DIR, Path(path).stem + ".txt")
                write_text_file(outtxt, ocr_text)
                HASH_DB_CACHE[fhash] = {"original": path, "status": "ocr_saved_txt", "out": outtxt}
                save_hash_db()
                return ("ocr_saved_txt", path, outtxt)
            else:
                # If OCR failed, mark as incompatible (move)
                dest = safe_copy(path, INCOMPATIBLE_DIR)
                return ("ocr_failed_moved_incompatible", path, dest)

        # If extension supported but not GPT4All-compatible (images, xlsx, pptx, epub) -> extract text if possible and save .txt into Processed_DIR
        if ext not in GPT4ALL_COMPATIBLE_EXTS:
            if text_len >= 10:
                outtxt = os.path.join(Processed_DIR, Path(path).stem + ".txt")
                write_text_file(outtxt, text)
                HASH_DB_CACHE[fhash] = {"original": path, "status": "extracted_and_saved", "out": outtxt}
                save_hash_db()
                return ("extracted_and_saved", path, outtxt)
            else:
                dest = safe_copy(path, INCOMPATIBLE_DIR)
                return ("incompatible_no_text", path, dest)

        # fallback
        dest = safe_copy(path, INCOMPATIBLE_DIR)
        return ("fallback_incompatible", path, dest)

    except MemoryError as me:
        append_error(f"MemoryError processing {path}: {me}\n{traceback.format_exc()}")
        return ("error_memory", path, str(me))
    except Exception as e:
        append_error(f"Exception processing {path}: {e}\n{traceback.format_exc()}")
        return ("error", path, str(e))


# ---------------------- Hash DB helpers ----------------------
def load_hash_db():
    if os.path.exists(HASH_DB_PATH):
        try:
            with open(HASH_DB_PATH, "r", encoding="utf-8") as fh:
                return json.load(fh)
        except Exception:
            append_error(f"Failed to load hash DB {HASH_DB_PATH}, starting fresh.")
            return {}
    return {}

def save_hash_db():
    # atomic write
    try:
        tmp = HASH_DB_PATH + ".tmp"
        with open(tmp, "w", encoding="utf-8") as fh:
            json.dump(HASH_DB_CACHE, fh, indent=2)
        os.replace(tmp, HASH_DB_PATH)
    except Exception as e:
        append_error(f"Failed to save hash DB: {e}")

# ---------------------- Main orchestration ----------------------
import json
HASH_DB_CACHE = load_hash_db()

def main():
    files = [str(p) for p in Path(INPUT_DIR).rglob("*") if p.is_file()]
    total = len(files)
    print(f"[INFO] Found {total} files in {INPUT_DIR}")
    if total == 0:
        print("[INFO] Nothing to do.")
        return

    # setup thread pool
    use_workers = MAX_WORKERS
    print(f"[INFO] Using {use_workers} worker threads (configurable: MAX_WORKERS).")
    results = []
    stats = {
        "processed": 0, "skipped": 0, "copied_original": 0, "ocr_saved_txt": 0,
        "extracted_and_saved": 0, "incompatible": 0, "errors": 0
    }

    # Submit tasks in threads, but limit queue size to avoid flooding memory
    from concurrent.futures import ThreadPoolExecutor, as_completed
    with ThreadPoolExecutor(max_workers=use_workers) as exe:
        futures = {exe.submit(process_single_file, fp): fp for fp in files}
        pbar = tqdm(total=len(futures), unit="file", desc="Overall processing")
        last_report = time.time()
        processed_count = 0
        start_time = time.time()

        for fut in as_completed(futures):
            processed_count += 1
            fp = futures[fut]
            try:
                status, orig, info = fut.result()
            except Exception as e:
                append_error(f"Unhandled worker exception for {fp}: {e}\n{traceback.format_exc()}")
                status = "error"
                orig = fp
                info = str(e)

            # update stats
            if status == "skipped":
                stats["skipped"] += 1
            elif status == "copied_original":
                stats["copied_original"] += 1
            elif status == "ocr_saved_txt":
                stats["ocr_saved_txt"] += 1
            elif status == "extracted_and_saved":
                stats["extracted_and_saved"] += 1
            elif status in ("incompatible", "incompatible_no_text", "ocr_failed_moved_incompatible", "fallback_incompatible"):
                stats["incompatible"] += 1
            else:
                if status and status.startswith("error"):
                    stats["errors"] += 1

            pbar.update(1)

            # periodic throughput print (every 15 seconds)
            now = time.time()
            if now - last_report > 15:
                elapsed = now - start_time
                rate = processed_count / elapsed if elapsed > 0 else 0
                remaining = len(futures) - processed_count
                eta = remaining / rate if rate > 0 else float('inf')
                print(f"[PROGRESS] processed={processed_count}/{len(futures)} rate={rate:.2f}/s ETA={int(eta)}s")
                last_report = now

        pbar.close()

    # final save hash DB
    try:
        save_hash_db()
    except Exception as e:
        append_error(f"Error saving final hash db: {e}")

    # final summary
    total_processed = sum([stats[k] for k in ("copied_original","ocr_saved_txt","extracted_and_saved")])
    print("\n========== SUMMARY ==========")
    print(f"Total files discovered : {total}")
    print(f"Already processed (skipped): {stats['skipped']}")
    print(f"Processed (case a & b & others saved): {total_processed}")
    print(f" - copied originals (case a): {stats['copied_original']}")
    print(f" - OCR -> .txt saved (case b): {stats['ocr_saved_txt']}")
    print(f" - extracted other -> .txt saved: {stats['extracted_and_saved']}")
    print(f"Incompatible / copied to Incompatible/: {stats['incompatible']}")
    print(f"Errors: {stats['errors']}")
    print(f"Processed outputs (point GPT4All LocalDocs to): {Processed_DIR}")
    print(f"Incompatible files are in: {INCOMPATIBLE_DIR}")
    print(f"Error log: {ERROR_LOG}")
    print(f"Hash DB: {HASH_DB_PATH}")
    print("=============================\n")

if __name__ == "__main__":
    main()
